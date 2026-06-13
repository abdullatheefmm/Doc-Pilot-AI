from __future__ import annotations

import csv
import io
import json
import time as _time
import uuid
from pathlib import Path
from typing import Literal

import os

from fastapi import FastAPI, File, HTTPException, Query, UploadFile, Depends, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from auth import require_admin, get_current_user, require_active_user, verify_user_password
from supabase_client import supabase
from pydantic import BaseModel, Field
from pypdf import PdfReader

from chunking import normalize_whitespace
import analytics
import knowledge_domains
import rag_qa
import retrieval
from cache import answer_cache

app = FastAPI(title="DocPilot AI — GenAI-Powered Knowledge Management System")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

DATA_DIR = Path("data/fastapi_docs")
UPLOAD_DIR = DATA_DIR / "uploads"
SUPPORTED_UPLOADS = {
    ".pdf",   # PDF documents
    ".txt",   # Plain text
    ".md",    # Markdown
    ".docx",  # Microsoft Word
    ".doc",   # Legacy Word (treated as text)
    ".pptx",  # Microsoft PowerPoint
    ".ppt",   # Legacy PowerPoint (treated as text)
    ".xlsx",  # Microsoft Excel
    ".xls",   # Legacy Excel (treated as text)
    ".csv",   # Comma-separated values
    ".json",  # JSON data files
}
ADMIN_PASSWORD = os.getenv("ADMIN_PASSWORD", "admin123")


# ─── Request / Response Models ──────────────────────────────────────

class QueryRequest(BaseModel):
    query: str = Field(..., min_length=1)
    session_id: str | None = None
    user_id: str | None = None
    user_role: str = "general"
    top_k: int = 4
    threshold: float = 0.15
    retrieval_mode: str = "hybrid"
    domain: str | None = None
    model: str | None = None
    incognito: bool = False
    search_mode: str = "internal"


class SourceItem(BaseModel):
    document: str
    text: str


class RetrievalScore(BaseModel):
    document: str
    score: float


class TimingInfo(BaseModel):
    rewrite_ms: float = 0
    retrieve_ms: float = 0
    generate_ms: float = 0


class QueryResponse(BaseModel):
    answer: str
    sources: list[SourceItem]
    confidence: float
    session_id: str
    rewritten_query: str
    intent: str
    cached: bool = False
    response_time_ms: float = 0
    timing: TimingInfo = TimingInfo()
    retrieval_scores: list[RetrievalScore] = []
    domain: str = "all"


class FeedbackRequest(BaseModel):
    session_id: str
    helpful: bool
    message_id: str | None = None
    query: str | None = None
    answer: str | None = None


class DomainCreateRequest(BaseModel):
    id: str = Field(..., min_length=1)
    name: str = Field(..., min_length=1)
    color: str = "#64748b"
    icon: str = "📁"


class UrlIngestRequest(BaseModel):
    url: str = Field(..., min_length=1)
    password: str = Field(..., min_length=1)
    domain: str | None = None
    user_id: str | None = None


# ─── Helpers ────────────────────────────────────────────────────────

def ensure_directories() -> None:
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)


def _extract_pdf(file_path: Path) -> str:
    """Extract text from a PDF file using PyMuPDF and OCR for images."""
    try:
        from vision_extractor import process_pdf_with_vision
        return process_pdf_with_vision(file_path)
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"PDF extraction failed: {exc}") from exc


def _extract_docx(file_path: Path) -> str:
    """Extract text from a DOCX file, preserving headings and paragraphs."""
    try:
        import docx  # python-docx
        doc = docx.Document(str(file_path))
        lines: list[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            # Tag headings for better chunking context
            if para.style and para.style.name.startswith("Heading"):
                level = para.style.name.replace("Heading ", "")
                lines.append(f"{'#' * int(level) if level.isdigit() else '#'} {text}")
            else:
                lines.append(text)
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    lines.append(" | ".join(cells))
        return "\n".join(lines)
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"DOCX extraction failed: {exc}") from exc


def _extract_pptx(file_path: Path) -> str:
    """Extract text from a PPTX file, slide by slide."""
    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        lines: list[str] = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            lines.append(f"--- Slide {slide_num} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
        return "\n".join(lines)
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"PPTX extraction failed: {exc}") from exc


def _extract_xlsx(file_path: Path) -> str:
    """Extract text from an XLSX file, sheet by sheet."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(file_path), data_only=True)
        lines: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines.append(f"--- Sheet: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    lines.append(" | ".join(cells))
        return "\n".join(lines)
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"XLSX extraction failed: {exc}") from exc


def _extract_csv(file_path: Path) -> str:
    """Extract text from a CSV file."""
    try:
        content = file_path.read_text(encoding="utf-8", errors="ignore")
        reader = csv.reader(io.StringIO(content))
        rows = [" | ".join(cell.strip() for cell in row if cell.strip()) for row in reader]
        return "\n".join(r for r in rows if r)
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"CSV extraction failed: {exc}") from exc


def _extract_json(file_path: Path) -> str:
    """Extract readable text from a JSON file."""
    try:
        data = json.loads(file_path.read_text(encoding="utf-8", errors="ignore"))
        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"JSON extraction failed: {exc}") from exc


def read_uploaded_text(file_path: Path) -> str:
    """Dispatch to the correct extractor based on file extension."""
    suffix = file_path.suffix.lower()

    if suffix in {".txt", ".md", ".doc", ".ppt", ".xls"}:
        # Legacy binary formats served as plain text (best-effort)
        return file_path.read_text(encoding="utf-8", errors="ignore")

    if suffix == ".pdf":
        return _extract_pdf(file_path)

    if suffix == ".docx":
        return _extract_docx(file_path)

    if suffix == ".pptx":
        return _extract_pptx(file_path)

    if suffix == ".xlsx":
        return _extract_xlsx(file_path)

    if suffix == ".csv":
        return _extract_csv(file_path)

    if suffix == ".json":
        return _extract_json(file_path)

    raise HTTPException(status_code=400, detail=f"Unsupported file type: '{suffix}'.")


# ─── Lifecycle ──────────────────────────────────────────────────────

@app.on_event("startup")
def startup_event() -> None:
    ensure_directories()
    retrieval.refresh_index()
    rag_qa.init_feedback_store()
    knowledge_domains.init_domains_db()
    analytics.init_analytics_db()


# ─── Core Q&A ───────────────────────────────────────────────────────

@app.post("/api/ask", response_model=QueryResponse)
def ask_question(request: QueryRequest):
    session_id = request.session_id or str(uuid.uuid4())

    try:
        if request.search_mode == "web":
            import web_agent
            web_generator = web_agent.stream_web_answer(request.query)
            citations = next(web_generator)
            answer_text = "".join(list(web_generator))
            if not request.incognito:
                rag_qa.append_to_history(session_id, "user", request.query.strip(), user_id=request.user_id)
                rag_qa.append_to_history(session_id, "assistant", answer_text, user_id=request.user_id)
            return {
                "answer": answer_text,
                "sources": [{"document": c, "text": "Web Result"} for c in citations],
                "confidence": 0.8,
                "session_id": session_id,
                "rewritten_query": request.query,
                "intent": "web_search",
                "cached": False,
                "response_time_ms": 0,
                "timing": {"rewrite_ms": 0, "retrieve_ms": 0, "generate_ms": 0},
                "retrieval_scores": [],
                "domain": "web",
            }
            
        return rag_qa.generate_answer(
            query=request.query.strip(),
            session_id=session_id,
            top_k=request.top_k,
            threshold=request.threshold,
            mode=request.retrieval_mode,
            domain=request.domain,
            user_id=request.user_id or "",
            user_role=request.user_role,
            model=request.model,
            search_mode=request.search_mode,
        )
    except RuntimeError as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Unexpected backend error: {exc}") from exc


@app.post("/api/ask/stream")
def ask_question_stream(request: QueryRequest):
    """Stream answer tokens via Server-Sent Events."""
    session_id = request.session_id or str(uuid.uuid4())

    if request.search_mode == "web":
        import web_agent
        def web_stream():
            web_generator = web_agent.stream_web_answer(request.query)
            citations = next(web_generator)
            
            meta = {
                "type": "meta",
                "confidence": 0.8,
                "intent": "web_search",
                "sources": [{"document": c, "text": "Web Result"} for c in citations],
                "retrieval_scores": [],
                "rewritten_query": request.query,
                "session_id": session_id,
                "retrieve_ms": 0,
                "domain": "web",
            }
            yield f"data: {json.dumps(meta)}\n\n"
            
            t0 = _time.perf_counter()
            full_answer = []
            for token in web_generator:
                full_answer.append(token)
                yield f"data: {json.dumps({'type': 'token', 'content': token})}\n\n"
                
            generate_ms = (_time.perf_counter() - t0) * 1000
            
            answer_text = "".join(full_answer)
            if not request.incognito:
                rag_qa.append_to_history(session_id, "user", request.query.strip(), user_id=request.user_id)
                rag_qa.append_to_history(session_id, "assistant", answer_text, user_id=request.user_id)
                
            done = {"type": "done", "generate_ms": round(generate_ms, 1)}
            yield f"data: {json.dumps(done)}\n\n"
        return StreamingResponse(web_stream(), media_type="text/event-stream")

    # Internal RAG Mode
    history_window = rag_qa.build_history_window(session_id)
    rewritten_query = rag_qa.rewrite_query(request.query.strip(), history_window)
    results, retrieve_ms = retrieval.retrieve(
        rewritten_query,
        user_id=request.user_id,
        user_role=request.user_role,
        top_k=request.top_k,
        threshold=request.threshold,
        mode=request.retrieval_mode,
        domain=request.domain,
    )
    confidence = rag_qa.compute_confidence(results)
    intent = rag_qa.detect_intent(request.query)
    sources = rag_qa.format_sources(results)
    retrieval_scores = rag_qa.format_retrieval_scores(results)

    def event_stream():
        full_answer = []
        # Send metadata first
        meta = {
            "type": "meta",
            "confidence": confidence,
            "intent": intent,
            "sources": sources,
            "retrieval_scores": retrieval_scores,
            "rewritten_query": rewritten_query,
            "session_id": session_id,
            "retrieve_ms": retrieve_ms,
            "domain": request.domain or "all",
        }
        yield f"data: {json.dumps(meta)}\n\n"

        t0 = _time.perf_counter()
        for token in rag_qa.generate_answer_stream(
            request.query.strip(), rewritten_query, results, session_id
        ):
            full_answer.append(token)
            yield f"data: {json.dumps({'type': 'token', 'content': token})}\n\n"
        generate_ms = (_time.perf_counter() - t0) * 1000

        answer_text = "".join(full_answer)
        if not request.incognito:
            rag_qa.append_to_history(session_id, "user", request.query.strip(), user_id=request.user_id)
            rag_qa.append_to_history(session_id, "assistant", answer_text, user_id=request.user_id)

        done = {"type": "done", "generate_ms": round(generate_ms, 1)}
        yield f"data: {json.dumps(done)}\n\n"

        # Log analytics
        try:
            analytics.log_query(
                query=request.query.strip(),
                answer=answer_text,
                confidence=confidence,
                response_time_ms=round(retrieve_ms + generate_ms, 1),
                sources=[r["document"] for r in results],
                domain=request.domain,
                intent=intent,
                cached=False,
                session_id=session_id,
            )
        except Exception:
            pass

    return StreamingResponse(event_stream(), media_type="text/event-stream")


# ─── Document Management ───────────────────────────────────────────

@app.post("/api/upload")
async def upload_document(
    file: UploadFile = File(...),
    password: str = Query(..., min_length=1),
    domain: str = Query("general"),
    user: dict = Depends(require_active_user),
):
    verify_user_password(user["email"], password)

    ensure_directories()

    suffix = Path(file.filename or "").suffix.lower()
    if suffix not in SUPPORTED_UPLOADS:
        supported = ", ".join(sorted(SUPPORTED_UPLOADS))
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type '{suffix}'. Supported: {supported}"
        )

    upload_path = UPLOAD_DIR / file.filename
    upload_path.write_bytes(await file.read())

    extracted_text = normalize_whitespace(read_uploaded_text(upload_path))
    normalized_path = DATA_DIR / f"{upload_path.stem}.txt"
    normalized_path.write_text(extracted_text, encoding="utf-8")

    # Assign domain
    knowledge_domains.assign_document_domain(normalized_path.name, domain)

    # Log to audit_logs
    if supabase:
        try:
            supabase.table("audit_logs").insert({
                "action_type": "upload_document",
                "details": {"filename": normalized_path.name, "original_filename": file.filename, "domain": domain, "email": user["email"]}
            }).execute()
        except Exception as e:
            print(f"Failed to log upload audit: {e}")

    # Generate AI summary & Trust Score
    summary = ""
    try:
        from summarization import generate_summary, generate_trust_score
        import trust_metrics
        import graph_extraction
        summary = generate_summary(extracted_text)
        knowledge_domains.save_document_summary(normalized_path.name, summary)
        
        trust_score = generate_trust_score(extracted_text)
        trust_metrics.save_ai_score(normalized_path.name, trust_score)
        
        # Process graph entities dynamically
        user_display_name = user.get("full_name") or user.get("email")
        graph_extraction.process_document_for_graph(normalized_path.name, extracted_text, domain, user_display_name)
    except Exception as e:
        print(f"Summary/Trust/Graph generation failed: {e}")

    retrieval.refresh_index()
    answer_cache.invalidate_all()

    return {
        "filename": normalized_path.name,
        "status": "Indexed successfully",
        "documents_indexed": len(retrieval.get_documents()),
        "domain": domain,
        "summary": summary,
    }


@app.get("/api/documents")
async def get_documents():
    docs = retrieval.get_document_inventory()
    domains = knowledge_domains.get_all_document_domains()
    summaries = knowledge_domains.get_all_summaries()
    
    import analytics
    import trust_metrics
    access_counts = analytics.get_document_access_counts()
    trust_metrics_map = trust_metrics.get_all_trust_metrics()
    
    upload_metadata = {}
    if supabase:
        try:
            res = supabase.table("audit_logs").select("created_at, user_id, details").eq("action_type", "upload_document").order("created_at", desc=False).execute()
            for log in res.data:
                details = log.get("details", {})
                if isinstance(details, str):
                    try:
                        import json
                        details = json.loads(details)
                    except:
                        details = {}
                filename = details.get("filename") if isinstance(details, dict) else None
                if filename:
                    # Overwrite so we keep the oldest or newest? newest is fine
                    upload_metadata[filename] = {
                        "uploaded_by": details.get("email", log.get("user_id", "Unknown")),
                        "date": log.get("created_at", "Unknown").split("T")[0]
                    }
        except Exception as e:
            print(f"Failed to fetch upload metadata: {e}")

    for doc in docs:
        doc_name = doc["name"]
        doc["domain"] = domains.get(doc_name, "general")
        doc["summary"] = summaries.get(doc_name, "")
        meta = upload_metadata.get(doc_name, {})
        doc["uploaded_by"] = meta.get("uploaded_by", "Unknown")
        doc["date"] = meta.get("date", "Unknown")
        doc["access_count"] = access_counts.get(doc_name, 0)
        doc["trust_metrics"] = trust_metrics_map.get(doc_name, {"document_name": doc_name, "upvotes": 0, "downvotes": 0, "ai_score": 85})
        
    return {"documents": docs}


@app.post("/api/documents/{doc_name}/vote")
async def vote_document(doc_name: str, payload: dict):
    vote_type = payload.get("vote")
    if vote_type in ("up", "down"):
        import trust_metrics
        trust_metrics.vote_document(doc_name, vote_type)
    return {"status": "ok"}


@app.delete("/api/documents/{filename}")
def delete_document(filename: str, password: str = Query(..., min_length=1), user: dict = Depends(require_active_user)):
    verify_user_password(user["email"], password)

    file_path = DATA_DIR / filename
    if file_path.exists():
        try:
            file_path.unlink()
        except Exception:
            pass

    for f in UPLOAD_DIR.glob(f"{Path(filename).stem}.*"):
        try:
            f.unlink()
        except Exception:
            pass

    # Manually delete vectors from Supabase
    retrieval.delete_document_from_index(filename)
    knowledge_domains.delete_document_metadata(filename)
    
    import graph_extraction
    graph_extraction.remove_document_from_graph(filename)
    
    retrieval.refresh_index()
    answer_cache.invalidate_all()
    
    if supabase:
        try:
            supabase.table("audit_logs").insert({
                "action_type": "delete_document",
                "user_id": user["email"],
                "domain": "general",
                "details": {"filename": filename}
            }).execute()
        except Exception:
            pass
            
    return {"status": "deleted", "filename": filename}


@app.get("/api/documents/{doc_id}/summary")
def get_document_summary(doc_id: str):
    summary = knowledge_domains.get_document_summary(f"{doc_id}.txt")
    return {"document": doc_id, "summary": summary or "No summary available."}


# ─── URL Ingestion ──────────────────────────────────────────────────

@app.post("/api/ingest-url")
def ingest_url(request: UrlIngestRequest):
    if request.password != ADMIN_PASSWORD:
        raise HTTPException(status_code=403, detail="Invalid admin password.")

    try:
        from ingestion import ingest_url as _ingest, save_ingested
        filename, text = _ingest(request.url)
        save_ingested(filename, text)
    except Exception as exc:
        raise HTTPException(status_code=400, detail=f"Failed to ingest URL: {exc}") from exc

    # Assign domain
    knowledge_domains.assign_document_domain(filename, request.domain or "general")

    # Log to audit_logs
    if supabase:
        try:
            supabase.table("audit_logs").insert({
                "action_type": "upload_document",
                "details": {"filename": filename, "original_filename": request.url, "domain": request.domain or "general", "email": request.user_id or "Unknown"}
            }).execute()
        except Exception as e:
            print(f"Failed to log upload audit: {e}")

    # Generate AI summary
    try:
        from summarization import generate_summary
        summary = generate_summary(text)
        knowledge_domains.save_document_summary(filename, summary)
    except Exception:
        pass

    retrieval.refresh_index()
    answer_cache.invalidate_all()
    
    return {"status": "success", "filename": filename}

    retrieval.refresh_index()
    answer_cache.invalidate_all()

    return {
        "filename": filename,
        "status": "Indexed successfully",
        "documents_indexed": len(retrieval.get_documents()),
        "domain": request.domain or "general",
    }


# ─── Knowledge Domains ─────────────────────────────────────────────

@app.get("/api/domains")
async def list_domains():
    domains = knowledge_domains.list_domains()
    doc_domains = knowledge_domains.get_all_document_domains()
    for d in domains:
        d["document_count"] = sum(1 for v in doc_domains.values() if v == d["id"])
    return {"domains": domains}


@app.post("/api/domains")
def create_domain(request: DomainCreateRequest):
    domain = knowledge_domains.create_domain(request.id, request.name, request.color, request.icon)
    return {"domain": domain}


@app.put("/api/documents/{doc_name}/domain")
def assign_domain(doc_name: str, domain: str = Query(..., min_length=1)):
    knowledge_domains.assign_document_domain(doc_name, domain)
    return {"status": "assigned", "document": doc_name, "domain": domain}


# ─── Chat History ───────────────────────────────────────────────────

@app.get("/api/history/sessions")
def get_sessions(user_id: str | None = None):
    return {"sessions": rag_qa.get_all_sessions(user_id)}

@app.get("/api/history")
def get_history(session_id: str = Query(..., min_length=1)):
    return {"session_id": session_id, "history": rag_qa.get_history(session_id)}

@app.delete("/api/history/{session_id}")
def clear_history(session_id: str):
    rag_qa.clear_history(session_id)
    return {"status": "cleared"}


# ─── Feedback ───────────────────────────────────────────────────────

@app.post("/api/feedback")
def submit_feedback(request: FeedbackRequest):
    rag_qa.save_feedback(
        session_id=request.session_id,
        helpful=request.helpful,
        message_id=request.message_id,
        query=request.query,
        answer=request.answer,
    )
    return {"status": "saved"}


@app.get("/api/feedback")
def list_feedback():
    return {"feedback": rag_qa.export_feedback()}


# ─── Analytics ──────────────────────────────────────────────────────

@app.get("/api/analytics/dashboard")
def get_analytics_dashboard():
    data = analytics.get_dashboard_data()
    data["cache_stats"] = answer_cache.stats()
    return data


# ─── Export ─────────────────────────────────────────────────────────

@app.get("/api/export/chat/{session_id}")
def export_chat(session_id: str):
    history = rag_qa.get_history(session_id)
    if not history:
        return {"markdown": "# Chat Export\n\nNo messages in this session."}

    lines = ["# DocPilot AI — Chat Export", f"**Session:** {session_id}\n", "---\n"]
    for msg in history:
        role = "**You**" if msg["role"] == "user" else "**Assistant**"
        lines.append(f"{role}:\n{msg['content']}\n")
    return {"markdown": "\n".join(lines)}


@app.get("/api/export/knowledge-base")
def export_knowledge_base():
    docs = retrieval.get_document_inventory()
    summaries = knowledge_domains.get_all_summaries()
    domains_map = knowledge_domains.get_all_document_domains()

    lines = ["# DocPilot AI — Knowledge Base Export\n"]
    for doc in docs:
        lines.append(f"## {doc['name']}")
        lines.append(f"- **Size:** {doc['size']}")
        lines.append(f"- **Domain:** {domains_map.get(doc['name'], 'general')}")
        summary = summaries.get(doc["name"], "No summary available.")
        lines.append(f"- **Summary:** {summary}\n")
    return {"markdown": "\n".join(lines)}

@app.get("/api/admin/analytics")
def get_admin_analytics(user: dict = Depends(require_admin)):
    import analytics
    from datetime import datetime, timedelta
    from supabase_client import supabase
    
    if not supabase: return {"status": "error"}
    
    now = datetime.utcnow()
    last_week = (now - timedelta(days=7)).isoformat()
    
    logs_res = supabase.table("audit_logs").select("*").gte("created_at", last_week).execute()
    queries = [l for l in logs_res.data if l["action_type"] == "chat_query"]
    
    dates = {}
    for d in range(7):
        date_str = (now - timedelta(days=6-d)).strftime("%Y-%m-%d")
        dates[date_str] = 0
        
    for q in queries:
        day = q["created_at"].split("T")[0]
        if day in dates: dates[day] += 1
        
    trend = [{"date": k, "queries": v} for k, v in dates.items()]
    
    cost_res = supabase.table("audit_logs").select("details").eq("action_type", "chat_query").execute()
    total_tokens = sum((l.get("details") or {}).get("tokens_used", 0) for l in cost_res.data)
    est_cost = (total_tokens / 1000) * 0.002
    
    # Calculate knowledge gaps from low confidence queries
    low_conf = [q for q in queries if (q.get("details") or {}).get("confidence", 1.0) < 0.5]
    gap_counts = {}
    for q in low_conf:
        query_text = (q.get("details") or {}).get("query", "Unknown")
        gap_counts[query_text] = gap_counts.get(query_text, 0) + 1
    
    gap_summary = [{"query": k, "misses": v, "confidence": (1 - min(1.0, v*0.1))*100} for k, v in list(gap_counts.items())[:5]]
    
    # Calculate top documents
    doc_counts = {}
    for q in queries:
        docs = (q.get("details") or {}).get("retrieved_docs", [])
        for d in docs:
            doc_counts[d] = doc_counts.get(d, 0) + 1
            
    doc_summary = [{"doc": k, "hits": v} for k, v in sorted(doc_counts.items(), key=lambda x: x[1], reverse=True)[:5]]
    
    return {
        "trend": trend,
        "metrics": {
            "total_queries": len(queries),
            "estimated_cost": round(est_cost, 4),
            "total_tokens": total_tokens,
            "cache_hits": len([q for q in queries if (q.get("details") or {}).get("cache_hit")])
        },
        "knowledge_gaps": gap_summary,
        "top_documents": doc_summary
    }

@app.get("/api/knowledge-graph")
def get_knowledge_graph(domain: str = Query(None), view_type: str = Query(None), user: dict = Depends(require_active_user)):
    import graph_extraction
    user_role = (user.get("role") or "general").lower()
    is_super_admin = user_role == "super_admin"
    user_email = user.get("email") or ""
    
    target_domain = domain.lower() if domain else user_role
    return graph_extraction.get_dynamic_graph_for_user(target_domain, is_super_admin, user_email, view_type)

@app.get("/api/admin/audit_logs")
def get_audit_logs(user: dict = Depends(require_admin)):
    from supabase_client import supabase
    if not supabase: return {"status": "error"}
    res = supabase.table("audit_logs").select("*").order("created_at", desc=True).limit(500).execute()
    return {"logs": res.data}

@app.delete("/api/admin/audit_logs/purge")
def purge_audit_logs(days: int = 30, user: dict = Depends(require_admin)):
    from supabase_client import supabase
    from datetime import datetime, timedelta
    if not supabase: return {"status": "error"}
    cutoff_date = (datetime.utcnow() - timedelta(days=days)).isoformat()
    supabase.table("audit_logs").delete().lte("created_at", cutoff_date).execute()
    return {"status": "success"}

@app.delete("/api/admin/audit_logs/{log_id}")
def delete_audit_log(log_id: str, user: dict = Depends(require_admin)):
    from supabase_client import supabase
    if not supabase: return {"status": "error"}
    supabase.table("audit_logs").delete().eq("id", log_id).execute()
    return {"status": "success"}

@app.post("/api/admin/auto_suspend")
def toggle_auto_suspend(req: dict, user: dict = Depends(require_admin)):
    # Persist this setting in db or memory
    # For now returning success
    enabled = req.get("enabled", False)
    return {"status": "success", "enabled": enabled}

@app.get("/api/auth/resolve-id")
def resolve_auth_id(emp_id: str, is_admin_login: str = "false"):
    from supabase_client import supabase
    if not supabase:
        raise HTTPException(status_code=500, detail="Supabase client not configured")
        
    res = supabase.table("user_profiles").select("email, role").eq("employee_id", emp_id).execute()
    if not res.data:
        raise HTTPException(status_code=404, detail="Not Found")
        
    user_data = res.data[0]
    if is_admin_login == "true" and user_data.get("role") != "super_admin":
        raise HTTPException(status_code=403, detail="Not an admin")
        
    return {"email": user_data.get("email")}

@app.get("/api/admin/users")
def get_users(user: dict = Depends(require_admin)):
    from supabase_client import supabase
    if not supabase: return {"status": "error"}
    res = supabase.table("user_profiles").select("*").execute()
    return {"users": res.data}

@app.delete("/api/admin/users/{user_id}")
def delete_user(user_id: str, user: dict = Depends(require_admin)):
    from supabase_client import supabase
    if not supabase: return {"status": "error"}
    supabase.table("user_profiles").delete().eq("user_id", user_id).execute()
    return {"status": "success"}

from pydantic import BaseModel
class StatusUpdate(BaseModel):
    status: str

@app.put("/api/admin/users/{user_id}/status")
def update_user_status(user_id: str, status_data: StatusUpdate, user: dict = Depends(require_admin)):
    from supabase_client import supabase
    if not supabase: return {"status": "error"}
    
    if status_data.status not in ["active", "pending", "revoked"]:
        raise HTTPException(status_code=400, detail="Invalid status")
        
    supabase.table("user_profiles").update({"status": status_data.status}).eq("user_id", user_id).execute()
    return {"status": "success"}


@app.get("/api/admin/security_alerts")
def get_security_alerts(user: dict = Depends(require_admin)):
    from supabase_client import supabase
    if not supabase: return {"status": "error"}
    res = supabase.table("audit_logs").select("*").in_("action_type", ["unauthorized_access", "bulk_download", "suspicious_login"]).order("created_at", desc=True).limit(50).execute()
    return {"alerts": res.data}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True, access_log=False, log_level="info")
